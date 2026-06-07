"""Generate Custom Claude Skills Guide — DETAILED step-by-step (15 slides)."""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colour palette — deep indigo / violet
# ---------------------------------------------------------------------------
INDIGO  = RGBColor(0x1E, 0x1B, 0x4B)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LAVEND  = RGBColor(0xF5, 0xF3, 0xFF)
GREY    = RGBColor(0x6B, 0x72, 0x80)
ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
LVIOLET = RGBColor(0xC4, 0xB5, 0xFD)
DVIOLET = RGBColor(0x4C, 0x1D, 0x95)
GREEN   = RGBColor(0x05, 0x96, 0x69)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(r, g, b): return RGBColor(r, g, b)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_solid(shape, colour):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour

def add_rect(slide, left, top, width, height, colour):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    fill_solid(shape, colour)
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, word_wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.text_frame.word_wrap = word_wrap
    return txb

def clear_first_para(txb):
    txf = txb.text_frame
    txf.paragraphs[0].clear()
    return txf

def para(txf, text, size, bold=False, colour=WHITE, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = txf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return p

def add_slide_number(slide, num, total):
    txb = add_textbox(slide, 12.5, 7.1, 0.7, 0.3)
    txf = clear_first_para(txb)
    para(txf, f"{num}/{total}", 9, colour=GREY, align=PP_ALIGN.RIGHT)

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.1, INDIGO)
    txb = add_textbox(slide, 0.4, 0.08, 12.5, 0.55)
    txf = clear_first_para(txb)
    para(txf, title, 22, bold=True, colour=WHITE)
    if subtitle:
        txb2 = add_textbox(slide, 0.4, 0.62, 12.5, 0.4)
        txf2 = clear_first_para(txb2)
        para(txf2, subtitle, 13, colour=LVIOLET)

def card(slide, left, top, width, height, heading, lines,
         head_colour=VIOLET, body_colour=LAVEND):
    h = 0.38
    add_rect(slide, left, top, width, h, head_colour)
    txh = add_textbox(slide, left + 0.1, top + 0.05, width - 0.2, h - 0.05)
    txf = clear_first_para(txh)
    para(txf, heading, 13, bold=True, colour=WHITE)
    add_rect(slide, left, top + h, width, height - h, body_colour)
    txb = add_textbox(slide, left + 0.12, top + h + 0.1, width - 0.24, height - h - 0.15)
    txf2 = clear_first_para(txb)
    for line in lines:
        para(txf2, line, 11, colour=INDIGO, space_before=1)

def step_badge(slide, step_num, label):
    add_rect(slide, 11.1, 1.15, 2.0, 0.52, VIOLET)
    txb = add_textbox(slide, 11.12, 1.17, 1.95, 0.48)
    txf = clear_first_para(txb)
    para(txf, f"Step {step_num}  ·  {label}", 10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

def code_box(slide, left, top, width, height, lines, title=None):
    add_rect(slide, left, top, width, height, INDIGO)
    y_start = top + 0.12
    if title:
        txb0 = add_textbox(slide, left + 0.2, y_start, width - 0.4, 0.32)
        txf0 = clear_first_para(txb0)
        para(txf0, title, 11, bold=True, colour=VIOLET)
        y_start += 0.32
    txb = add_textbox(slide, left + 0.2, y_start, width - 0.4,
                      height - (y_start - top) - 0.1)
    txf = clear_first_para(txb)
    for line in lines:
        if line.startswith("#"):
            c = GREY
        elif line.startswith("//") or line.startswith("--"):
            c = _rgb(0xA0, 0x90, 0xD0)
        else:
            c = LVIOLET
        para(txf, line, 11, colour=c, space_before=2)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def s01_title(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, INDIGO)
    add_rect(sl, 0, 0, 0.22, 7.5, VIOLET)

    txb = add_textbox(sl, 0.6, 1.35, 12.3, 1.2)
    txf = clear_first_para(txb)
    para(txf, "Custom Skills for Claude Code", 40, bold=True, colour=WHITE)

    txb2 = add_textbox(sl, 0.6, 2.65, 12.3, 0.65)
    txf2 = clear_first_para(txb2)
    para(txf2, "Detailed Step-by-Step Guide  ·  Build · Deploy · Activate", 22,
         colour=LVIOLET)

    add_rect(sl, 0.6, 3.5, 8.0, 0.04, VIOLET)

    txb3 = add_textbox(sl, 0.6, 3.72, 12.3, 1.6)
    txf3 = clear_first_para(txb3)
    for line in [
        "15 slides  ·  Prerequisites → Folder → SKILL.md → Commit → Activate → Use",
        "",
        "Covers:   File structure  ·  Writing effective instructions",
        "          Committing & pushing  ·  Mobile app activation",
        "          Debugging  ·  Real examples from this repo",
    ]:
        para(txf3, line, 13, colour=_rgb(0xA5, 0x92, 0xF9), space_before=3)

    add_slide_number(sl, n, T)


def s02_what_are_skills(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "What Are Claude Skills?", "Custom slash commands backed by Markdown instructions")

    card(sl, 0.3, 1.25, 6.2, 5.85, "What is a Skill?",
         ["A custom slash command (e.g. /make-pptx)",
          "",
          "Defined by a single SKILL.md file in your repo",
          "",
          "Claude reads it at session start and registers",
          "the slash command automatically",
          "",
          "Lives at:  .claude/skills/<name>/SKILL.md",
          "",
          "No server, no plugin, no API key for setup",
          "",
          "Activated the moment you push to GitHub"],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 6.8, 1.25, 6.2, 5.85, "Why Build Skills?",
         ["Automate repetitive tasks in one command",
          "",
          "Share workflows with your team via Git —",
          "one push, everyone has the skill",
          "",
          "Works across Claude Code web, mobile + IDE",
          "",
          "Skills are version-controlled with your code",
          "",
          "Pass arguments just like a terminal command:",
          "  /clinical-query diabetic kidney disease",
          "",
          "Consistent, repeatable, documented behaviour"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, n, T)


def s03_prerequisites(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Prerequisites", "Three things you need before your first skill")

    prereqs = [
        ("GitHub Repository",
         ["✓  Public or private repo on GitHub",
          "",
          "✓  You have push access to the branch",
          "    your session will read",
          "",
          "✓  Skills must be committed — they are",
          "    just files in the repo, nothing more",
          "",
          "✓  Skills on a feature branch are only",
          "    active in sessions using that branch"]),
        ("Claude Mobile App",
         ["✓  Latest Claude app installed",
          "    (iOS or Android)",
          "",
          "✓  Signed in to your Anthropic account",
          "",
          "✓  Navigate to the Code tab at the",
          "    bottom of the screen",
          "",
          "✓  New sessions re-read the repo —",
          "    skills pushed after start appear next"]),
        ("Claude Code Session",
         ["✓  A web or desktop session connected",
          "    to your GitHub repository",
          "",
          "✓  Session reads .claude/ on startup",
          "",
          "✓  Skills registered = slash commands",
          "    appear in the / autocomplete list",
          "",
          "✓  To pick up new skills: start a fresh",
          "    session after pushing to GitHub"]),
    ]

    for i, (heading, lines) in enumerate(prereqs):
        lft = 0.3 + i * 4.35
        card(sl, lft, 1.25, 4.1, 5.85, heading, lines,
             head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, n, T)


def s04_folder(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 1 — Create the Folder Structure",
               "The folder name IS the slash command name")
    step_badge(sl, 1, "FOLDER")

    add_rect(sl, 0.3, 1.25, 7.8, 3.85, INDIGO)
    txb = add_textbox(sl, 0.55, 1.38, 7.3, 3.62)
    txf = clear_first_para(txb)
    para(txf, "Directory tree:", 12, bold=True, colour=LVIOLET)
    for line in [
        "your-repo/",
        "└── .claude/",
        "    └── skills/",
        "        ├── make-pptx/",
        "        │   ├── SKILL.md       ← required",
        "        │   └── pptxgenjs.md   ← optional reference",
        "        ├── clinical-query/",
        "        │   └── SKILL.md       ← required",
        "        └── my-new-skill/",
        "            └── SKILL.md       ← required",
    ]:
        para(txf, line, 13, colour=_rgb(0xC4, 0xB5, 0xFD), space_before=3)

    code_box(sl, 0.3, 5.32, 7.8, 1.88,
        ["# Create the folder (replace 'my-skill' with your command name)",
         "mkdir -p .claude/skills/my-skill",
         "",
         "# Then create the required SKILL.md",
         "touch .claude/skills/my-skill/SKILL.md"],
        title="Terminal commands")

    add_rect(sl, 8.4, 1.25, 4.6, 2.2, INDIGO)
    txb2 = add_textbox(sl, 8.6, 1.35, 4.2, 2.0)
    txf2 = clear_first_para(txb2)
    para(txf2, "The Golden Rule", 14, bold=True, colour=VIOLET)
    para(txf2, "Folder name = slash command", 12, colour=WHITE, space_before=6)
    para(txf2, "", 6, colour=WHITE)
    para(txf2, "make-pptx      →  /make-pptx", 12, colour=LVIOLET, space_before=3)
    para(txf2, "clinical-query →  /clinical-query", 12, colour=LVIOLET, space_before=3)
    para(txf2, "my-skill       →  /my-skill", 12, colour=LVIOLET, space_before=3)

    card(sl, 8.4, 3.65, 4.6, 2.2, "Naming Rules",
         ["Use lowercase and hyphens only",
          "No spaces, no underscores",
          "Keep it short — it's typed as /command",
          "Describe the action: deploy-to-aws",
          "  not: my-great-automation-tool"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_rect(sl, 8.4, 6.05, 4.6, 1.15, _rgb(0x2D, 0x27, 0x6B))
    txb3 = add_textbox(sl, 8.6, 6.13, 4.2, 0.98)
    txf3 = clear_first_para(txb3)
    para(txf3, "Only SKILL.md is required.", 12, bold=True, colour=VIOLET)
    para(txf3, "Add other .md files for longer reference guides — "
         "keep the folder tidy.", 11, colour=WHITE, space_before=4)

    add_slide_number(sl, n, T)


def s05_skillmd_header(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 2a — Write SKILL.md: Header",
               "The # heading must exactly match the folder name")
    step_badge(sl, "2a", "HEADER")

    add_rect(sl, 0.3, 1.25, 7.6, 5.95, INDIGO)
    txb = add_textbox(sl, 0.5, 1.35, 7.2, 5.75)
    txf = clear_first_para(txb)
    para(txf, "Full SKILL.md template", 12, bold=True, colour=LVIOLET)
    sections = [
        ("# /my-skill",                              VIOLET),
        ("",                                          WHITE),
        ("One-line description of what this does.",   WHITE),
        ("",                                          WHITE),
        ("## Usage",                                  _rgb(0x86, 0xEF, 0xAC)),
        ("/my-skill [optional-argument]",             WHITE),
        ("",                                          WHITE),
        ("## What this skill does",                   _rgb(0x86, 0xEF, 0xAC)),
        ("1. First action Claude will take",          WHITE),
        ("2. Second action",                          WHITE),
        ("3. Third action",                           WHITE),
        ("",                                          WHITE),
        ("## Instructions",                           _rgb(0x86, 0xEF, 0xAC)),
        ("When the user runs /my-skill:",             WHITE),
        ("- Do X (be specific: name the file/cmd)",   LVIOLET),
        ("- Then do Y",                               LVIOLET),
        ("- Finally deliver Z to the user",           LVIOLET),
        ("",                                          WHITE),
        ("## Output format",                          _rgb(0x86, 0xEF, 0xAC)),
        ("Describe exactly what Claude should return.",WHITE),
    ]
    for text, colour in sections:
        para(txf, text, 11, colour=colour, space_before=2)

    card(sl, 8.2, 1.25, 4.85, 1.55, "# heading — critical rule",
         ["Must start with  # /skill-name",
          "Must match the folder name exactly",
          "make-pptx/ folder → # /make-pptx",
          "Claude uses this to register the command"],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 8.2, 3.0, 4.85, 1.55, "Usage section",
         ["One line showing the command + args",
          "Use [brackets] for optional args",
          "Example: /my-skill [filename] [--dry-run]",
          "Keep it short — just the signature"],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 8.2, 4.75, 4.85, 2.45, "Output format section",
         ["Tell Claude exactly what to return:",
          "  A .pptx file delivered to user",
          "  A table printed in the chat",
          "  A commit pushed to the branch",
          "",
          "Without this section Claude decides",
          "the output format on its own"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, n, T)


def s06_skillmd_instructions(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 2b — Write SKILL.md: Instructions Section",
               "The most important part — Claude follows this literally")
    step_badge(sl, "2b", "INSTRUCTIONS")

    add_rect(sl, 0.3, 1.25, 12.7, 0.38, INDIGO)
    txh = add_textbox(sl, 0.5, 1.28, 12.3, 0.3)
    txfh = clear_first_para(txh)
    para(txfh, "Weak instructions vs. strong instructions", 12, bold=True, colour=WHITE)

    comparisons = [
        ("WEAK (vague)",
         "- Analyse the code\n- Make it better\n- Show the result",
         "TOO VAGUE",
         "Strong (specific)",
         "- Run: python -m markitdown output.pptx\n"
         "- Check for missing content, typos, wrong order\n"
         "- Convert to images: scripts/office/soffice.py --headless …\n"
         "- Report each issue found, including minor ones",
         "CORRECT"),
    ]

    add_rect(sl, 0.3, 1.83, 6.0, 0.35, _rgb(0x7C, 0x28, 0x00))
    txw = add_textbox(sl, 0.5, 1.86, 5.8, 0.28)
    txfw = clear_first_para(txw)
    para(txfw, "WEAK — do not write like this", 11, bold=True, colour=_rgb(0xFF, 0xAA, 0x66))

    code_box(sl, 0.3, 2.18, 6.0, 1.88,
        ["## Instructions",
         "When the user invokes /my-skill:",
         "- Analyse the code",
         "- Make it better",
         "- Show the result"])

    add_rect(sl, 6.7, 1.83, 6.3, 0.35, _rgb(0x1A, 0x5C, 0x2E))
    txs = add_textbox(sl, 6.9, 1.86, 6.1, 0.28)
    txfs = clear_first_para(txs)
    para(txfs, "STRONG — write like this", 11, bold=True, colour=_rgb(0x88, 0xFF, 0xBB))

    code_box(sl, 6.7, 2.18, 6.3, 1.88,
        ["## Instructions",
         "When the user runs /my-skill:",
         "1. Run:  python -m markitdown output.pptx",
         "2. Check for: missing content, typos, wrong order",
         "3. Fix each issue found",
         "4. Deliver output.pptx to the user",
         "5. Commit and push: git add output.pptx"])

    add_rect(sl, 0.3, 4.27, 12.7, 0.38, DVIOLET)
    txhb = add_textbox(sl, 0.5, 4.3, 12.3, 0.3)
    txfhb = clear_first_para(txhb)
    para(txfhb, "5 rules for effective instructions", 12, bold=True, colour=WHITE)

    rules = [
        ("Name files & commands",  "Say run scripts/make_pptx.py — not 'run the script'"),
        ("Specify what to deliver","Say 'deliver output.pptx to the user' — Claude won't assume"),
        ("Break into numbered steps","Claude follows numbered lists reliably and in order"),
        ("Use if/else for conditionals","If no args given, do X. If --dry-run passed, do Y instead."),
        ("Reference other .md files", "For long guides: 'Read editing.md for full details'"),
    ]

    for i, (rule, example) in enumerate(rules):
        bg = LAVEND if i % 2 == 0 else WHITE
        add_rect(sl, 0.3, 4.65 + i * 0.52, 12.7, 0.52, bg)
        txr = add_textbox(sl, 0.5, 4.7 + i * 0.52, 3.6, 0.42)
        txfr = clear_first_para(txr)
        para(txfr, rule, 11, bold=True, colour=DVIOLET)
        txe = add_textbox(sl, 4.3, 4.7 + i * 0.52, 8.5, 0.42)
        txfe = clear_first_para(txe)
        para(txfe, example, 11, colour=INDIGO)

    add_slide_number(sl, n, T)


def s07_supporting_files(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 3 — Add Supporting Reference Files",
               "Keep SKILL.md short — offload long docs to companion .md files")
    step_badge(sl, 3, "DOCS")

    card(sl, 0.3, 1.25, 6.2, 2.65, "When to add companion files",
         ["SKILL.md should stay short (< 100 lines)",
          "  so Claude reads it fully and reliably",
          "",
          "Move long reference content to other .md files",
          "  in the same skills folder",
          "",
          "Reference them from Instructions:",
          '  "Read pptxgenjs.md for full details"',
          '  "See editing.md for the editing workflow"'],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 6.7, 1.25, 6.3, 2.65, "What goes in companion files",
         ["Long code templates or boilerplate",
          "Step-by-step sub-guides (editing, QA, etc.)",
          "Colour palettes, design systems, style guides",
          "API reference snippets",
          "Troubleshooting checklists",
          "",
          "Claude will read these when SKILL.md",
          "explicitly tells it to"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_rect(sl, 0.3, 4.12, 12.7, 0.38, INDIGO)
    txh = add_textbox(sl, 0.5, 4.15, 12.3, 0.3)
    txfh = clear_first_para(txh)
    para(txfh, "Example: make-pptx skill file layout", 12, bold=True, colour=WHITE)

    code_box(sl, 0.3, 4.5, 6.0, 2.7,
        [".claude/skills/make-pptx/",
         "├── SKILL.md          ← short, ~50 lines",
         "│     Trigger, Usage, Instructions,",
         "│     Output format, QA checklist ref",
         "├── pptxgenjs.md      ← 200+ lines",
         "│     Full PptxGenJS API reference",
         "├── editing.md        ← 150+ lines",
         "│     Unpack/edit/repack workflow",
         "└── README.md         ← optional overview"])

    code_box(sl, 6.6, 4.5, 6.5, 2.7,
        ["# In SKILL.md — how to reference companion files:",
         "",
         "## Instructions",
         "When the user runs /make-pptx:",
         "1. Run scripts/make_pptx.py",
         "2. Run QA (see QA section below)",
         "3. For editing an existing file: Read editing.md",
         "4. For creating from scratch: Read pptxgenjs.md",
         "5. Deliver the .pptx file to the user"])

    add_slide_number(sl, n, T)


def s08_commit_push(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 4 — Commit & Push to GitHub",
               "Skills only activate once committed to the branch Claude reads")
    step_badge(sl, 4, "GIT PUSH")

    code_box(sl, 0.3, 1.25, 12.7, 3.05,
        ["# Stage just the skill files (be precise — don't use git add -A)",
         "git add .claude/skills/my-skill/",
         "",
         "# Commit with a descriptive message",
         'git commit -m "Add /my-skill custom skill"',
         "",
         "# Push to the branch your Claude Code session is reading",
         "git push -u origin main",
         "",
         "# If working on a feature branch:",
         "git push -u origin feature/my-branch",
         "# Note: skills on a feature branch are only active in sessions",
         "#       connected to that branch — merge to main for everyone"],
        title="Git commands")

    add_rect(sl, 0.3, 4.52, 12.7, 1.55, _rgb(0x5C, 0x20, 0x00))
    txw = add_textbox(sl, 0.5, 4.6, 12.3, 1.38)
    txfw = clear_first_para(txw)
    para(txfw, "Warning — branch scope", 14, bold=True, colour=ORANGE)
    para(txfw, ("Skills pushed to feature/my-branch are only visible in sessions connected to that branch. "
                "Merge to main (or whichever branch your team uses) for the skill to be universally available."),
         12, colour=WHITE, space_before=4)

    card(sl, 0.3, 6.25, 6.0, 1.0, "Check what was committed",
         ["git log --oneline -3",
          "git show --stat HEAD"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_rect(sl, 6.7, 6.25, 6.3, 1.0, _rgb(0x2D, 0x27, 0x6B))
    txb = add_textbox(sl, 6.9, 6.33, 6.0, 0.82)
    txf = clear_first_para(txb)
    para(txf, "Tip:", 12, bold=True, colour=VIOLET)
    para(txf, ("Claude re-reads .claude/ on every new session start. "
               "Skills pushed after session start appear in the next session only."),
         11, colour=LVIOLET, space_before=3)

    add_slide_number(sl, n, T)


def s09_open_app(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 5 — Open Claude App & Navigate to Code Tab",
               "Mobile: iOS or Android  ·  Desktop: Mac or Windows")
    step_badge(sl, 5, "OPEN APP")

    steps = [
        ("1", "Open Claude App",
         "iOS or Android.\nTap the Claude icon\nto open."),
        ("2", "Tap Code Tab",
         "Bottom navigation bar.\nSelect the 'Code' tab\n(laptop icon)."),
        ("3", "Tap New Session",
         "Or continue an existing\nsession that's already\nconnected to your repo."),
        ("4", "Select Repository",
         "Choose the GitHub repo\nthat has your\n.claude/skills/ folder."),
        ("5", "Session Starts",
         "Claude clones the repo\nand scans .claude/\non startup."),
    ]

    box_w, box_h, gap = 2.28, 3.5, 0.27
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    start = (13.33 - total_w) / 2

    for i, (num, title, body) in enumerate(steps):
        lft = start + i * (box_w + gap)
        add_rect(sl, lft, 1.3, box_w, 0.5, DVIOLET)
        txn = add_textbox(sl, lft, 1.33, box_w, 0.44)
        txfn = clear_first_para(txn)
        para(txfn, f"Step {num}", 12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, lft, 1.8, box_w, box_h - 0.5, INDIGO)
        txh = add_textbox(sl, lft + 0.12, 1.9, box_w - 0.24, 0.42)
        txfh = clear_first_para(txh)
        para(txfh, title, 13, bold=True, colour=VIOLET)
        txb = add_textbox(sl, lft + 0.12, 2.38, box_w - 0.24, box_h - 1.1)
        txfb = clear_first_para(txb)
        for part in body.split("\n"):
            para(txfb, part, 11, colour=_rgb(0xCC, 0xBB, 0xFF), space_before=4)
        if i < len(steps) - 1:
            arr = add_textbox(sl, lft + box_w + 0.03, 2.9, gap, 0.38)
            txfa = clear_first_para(arr)
            para(txfa, "→", 16, bold=True, colour=VIOLET, align=PP_ALIGN.CENTER)

    add_rect(sl, 0.3, 5.02, 12.7, 1.08, _rgb(0x2D, 0x27, 0x6B))
    txb = add_textbox(sl, 0.5, 5.1, 12.3, 0.92)
    txf = clear_first_para(txb)
    para(txf, ("Skills are registered at session start. If you pushed a new skill after starting the current session, "
               "you must start a NEW session for the skill to appear in the / autocomplete list."),
         13, colour=LVIOLET)

    add_rect(sl, 0.3, 6.28, 12.7, 0.95, _rgb(0x1A, 0x5C, 0x2E))
    txb2 = add_textbox(sl, 0.5, 6.36, 12.3, 0.78)
    txf2 = clear_first_para(txb2)
    para(txf2, ("Verification:  In the new session, type  /  in the chat input and look for your skill name "
                "in the autocomplete dropdown. If it's there, you're ready."),
         12, colour=_rgb(0x88, 0xFF, 0xBB))

    add_slide_number(sl, n, T)


def s10_invoke(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 6 — Invoke Your Skill",
               "Type / to see registered skills and trigger them with arguments")
    step_badge(sl, 6, "USE SKILL")

    card(sl, 0.3, 1.25, 5.8, 5.85, "How to Trigger",
         ["1.  Type  /  in the chat input box",
          "",
          "2.  Autocomplete shows all registered skills",
          "",
          "3.  Select from the list or type the full name:",
          "       /make-pptx",
          "       /clinical-query",
          "",
          "4.  Optionally add arguments after the command:",
          "       /clinical-query diabetic kidney disease",
          "       /make-pptx --dry-run",
          "",
          "5.  Claude reads SKILL.md and executes",
          "    the Instructions section step-by-step",
          "",
          "6.  Arguments appear as args in the skill",
          "    and Claude uses them as context"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_rect(sl, 6.4, 1.25, 6.6, 5.85, INDIGO)
    txbh = add_textbox(sl, 6.6, 1.35, 6.2, 0.38)
    txfh = clear_first_para(txbh)
    para(txfh, "Live session example", 12, bold=True, colour=LVIOLET)

    examples = [
        ("You: /clinical-query",                      VIOLET,                   True),
        ("     diabetic patient with kidney",          WHITE,                    False),
        ("     complications",                         WHITE,                    False),
        ("",                                           WHITE,                    False),
        ("Claude: Reading SKILL.md…",                  GREY,                     False),
        ("        Running retrieval pipeline…",         GREY,                     False),
        ("",                                           WHITE,                    False),
        ("Rank  Score   Specialty",                    _rgb(0x86, 0xEF, 0xAC),  True),
        ("  1   0.985   Nephrology",                   WHITE,                    False),
        ("  2   0.702   Endocrinology",                WHITE,                    False),
        ("  3   0.611   Internal Medicine",             WHITE,                    False),
        ("",                                           WHITE,                    False),
        ("Clinical Summary:",                          _rgb(0x86, 0xEF, 0xAC),  True),
        ("The top result highlights stage-3 CKD",     LVIOLET,                  False),
        ("in a T2DM patient (GFR 38). Red flag:",     LVIOLET,                  False),
        ("proteinuria present across 3 of 5 notes.",  LVIOLET,                  False),
    ]

    txb2 = add_textbox(sl, 6.6, 1.82, 6.2, 5.1)
    txf2 = clear_first_para(txb2)
    for text, colour, bold in examples:
        para(txf2, text, 11, colour=colour, bold=bold, space_before=1)

    add_slide_number(sl, n, T)


def s11_passing_args(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 7 — Passing Arguments to Skills",
               "Skills accept arguments just like terminal commands")
    step_badge(sl, 7, "ARGS")

    add_rect(sl, 0.3, 1.25, 12.7, 0.38, INDIGO)
    txh = add_textbox(sl, 0.5, 1.28, 12.3, 0.3)
    txfh = clear_first_para(txh)
    para(txfh, "Syntax: /skill-name [arg1] [arg2] [--flag]", 12, bold=True, colour=WHITE)

    examples = [
        ("/make-pptx",
         "(no args)",
         "Run default: regenerate Clinical_Note_Retrieval_System.pptx"),
        ("/make-pptx more detailed step by step ppt for both",
         "free-form args",
         "Claude receives 'more detailed step by step ppt for both' as context"),
        ("/clinical-query diabetic kidney disease",
         "query string",
         "Runs retrieval with 'diabetic kidney disease' as the search query"),
        ("/clinical-query COPD with oxygen desaturation --top_k 10",
         "query + flag",
         "Claude extracts --top_k 10 and passes it to the retrieval call"),
    ]

    for i, (cmd, arg_type, result) in enumerate(examples):
        bg = LAVEND if i % 2 == 0 else WHITE
        add_rect(sl, 0.3, 1.83 + i * 0.68, 12.7, 0.68, bg)
        txc = add_textbox(sl, 0.5, 1.88 + i * 0.68, 4.8, 0.58)
        txfc = clear_first_para(txc)
        para(txfc, cmd, 11, bold=True, colour=DVIOLET)
        txat = add_textbox(sl, 5.5, 1.88 + i * 0.68, 1.8, 0.58)
        txfat = clear_first_para(txat)
        para(txfat, arg_type, 10, colour=GREY, italic=True)
        txr = add_textbox(sl, 7.5, 1.88 + i * 0.68, 5.3, 0.58)
        txfr = clear_first_para(txr)
        para(txfr, result, 11, colour=INDIGO)

    code_box(sl, 0.3, 4.6, 6.2, 2.6,
        ["## Instructions",
         "ARGUMENTS: {args}",
         "",
         "When the user runs /my-skill:",
         "- If no args: run the default action",
         "- If args contain '--dry-run': show what",
         "  would happen without doing it",
         "- If args contain a filename: use that file",
         "  instead of the default"],
        title="How to handle args in SKILL.md")

    card(sl, 6.7, 4.6, 6.3, 2.6, "Args best practices",
         ["Free-form args: Claude interprets context",
          "  naturally — no special parsing needed",
          "",
          "--flags are extracted by Claude naturally:",
          "  '--top_k 5', '--dry-run', '--format csv'",
          "",
          "Reference ARGUMENTS in your Instructions",
          "  section so Claude knows they exist",
          "",
          "Keep it simple — Claude handles ambiguity"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, n, T)


def s12_real_example_cq(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Real Example — /clinical-query Skill",
               "The full SKILL.md for clinical note retrieval")
    step_badge(sl, "EX 1", "EXAMPLE")

    code_box(sl, 0.3, 1.25, 8.0, 6.0,
        ["# /clinical-query",
         "",
         "Run a semantic search over the indexed clinical notes.",
         "Returns top-k ranked results with cosine scores and",
         "an LLM-generated clinical summary.",
         "",
         "## Usage",
         "/clinical-query <natural language clinical scenario>",
         "",
         "## What this skill does",
         "1. Calls the retrieval pipeline with the query",
         "2. Prints ranked results table (rank, score, specialty)",
         "3. Shows the LLM clinical summary",
         "4. Reports score spread across top-k results",
         "",
         "## Instructions",
         "When the user runs /clinical-query:",
         "1. Extract the query from ARGUMENTS (everything after /clinical-query)",
         "2. Run: python scripts/demo.py \"<query>\"",
         "3. Parse and display the results as a Markdown table",
         "4. Print the clinical summary below the table",
         "5. Note the score spread (rank-1 score minus rank-k score)",
         "",
         "## Output format",
         "A Markdown table: Rank | Score | Specialty | Description",
         "Followed by the LLM clinical summary paragraph."],
        title=".claude/skills/clinical-query/SKILL.md")

    card(sl, 8.6, 1.25, 4.4, 3.0, "Key design decisions",
         ["Query extracted from ARGUMENTS",
          "  — no special parsing needed",
          "",
          "Uses scripts/demo.py — named",
          "  explicitly so Claude doesn't guess",
          "",
          "Output format specified exactly:",
          "  Markdown table + summary paragraph",
          "",
          "Score spread reported — helps",
          "  the user assess confidence"],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 8.6, 4.45, 4.4, 2.8, "Result in chat",
         ["Rank  Score   Specialty",
          "1     0.985   Nephrology",
          "2     0.702   Endocrinology",
          "3     0.611   Internal Medicine",
          "",
          "Clinical Summary:",
          "The top result highlights stage-3 CKD",
          "in a T2DM patient (GFR 38)…"],
         head_colour=_rgb(0x1A, 0x5C, 0x2E), body_colour=WHITE)

    add_slide_number(sl, n, T)


def s13_real_example_pptx(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Real Example — /make-pptx Skill",
               "The full SKILL.md for PPTX generation")
    step_badge(sl, "EX 2", "EXAMPLE")

    code_box(sl, 0.3, 1.25, 8.0, 6.0,
        ["# /make-pptx",
         "",
         "Generate, QA, and deliver a .pptx presentation.",
         "",
         "## Usage",
         "/make-pptx [instructions]",
         "",
         "## What this skill does",
         "1. Runs scripts/make_pptx.py (or follows instructions)",
         "2. Extracts text with markitdown for content QA",
         "3. Converts to images for visual QA",
         "4. Fixes issues found in QA",
         "5. Delivers the .pptx file to the user",
         "6. Commits and pushes to the repo",
         "",
         "## Instructions",
         "When the user runs /make-pptx:",
         "1. Read the ARGUMENTS — if they contain specific instructions,",
         "   follow those; otherwise regenerate the default presentation",
         "2. Install deps: pip install python-pptx markitdown --quiet",
         "3. Run: python3 scripts/make_pptx.py",
         "4. QA: python -m markitdown output.pptx",
         "5. Fix any missing content, placeholders, or layout issues",
         "6. Deliver the .pptx file to the user",
         "7. Commit and push: git add *.pptx && git commit -m '...'"],
        title=".claude/skills/make-pptx/SKILL.md (excerpt)")

    card(sl, 8.6, 1.25, 4.4, 2.95, "What makes this skill robust",
         ["ARGUMENTS checked first — skill",
          "  adapts to instructions dynamically",
          "",
          "Deps installed inline — no manual",
          "  setup step for the user",
          "",
          "QA step is mandatory — not optional",
          "",
          "Commit + push included — user",
          "  doesn't need to do it manually",
          "",
          "References pptxgenjs.md and",
          "  editing.md for long guides"],
         head_colour=DVIOLET, body_colour=WHITE)

    card(sl, 8.6, 4.4, 4.4, 2.85, "Supporting files",
         [".claude/skills/make-pptx/",
          "  SKILL.md         ← this file",
          "  pptxgenjs.md     ← create from scratch",
          "  editing.md       ← edit existing .pptx",
          "",
          "Claude reads them when SKILL.md says:",
          "  'Read pptxgenjs.md for full details'"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, n, T)


def s14_debug(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 8 — Debug If Your Skill Doesn't Appear",
               "Systematic checklist for common problems")
    step_badge(sl, 8, "DEBUG")

    issues = [
        ("Skill not in / list",
         "The most common cause: you're on a different branch than what you pushed to.",
         ["Check: git branch --show-current",
          "Fix: git push -u origin <your-current-branch>",
          "Or: merge the skill branch to main"]),
        ("Skill pushed but still missing",
         "Claude reads .claude/ only at session start — a running session won't pick up new skills.",
         ["Fix: start a brand-new session",
          "Verify: type / — the skill should appear in autocomplete",
          "Check: git log --oneline -3 to confirm the commit reached GitHub"]),
        ("Skill appears but Claude ignores the steps",
         "Instructions are too vague — Claude interprets loosely rather than follows exactly.",
         ["Fix: name files, commands, and outputs explicitly",
          "Replace 'run the script' with 'run: python scripts/make_pptx.py'",
          "Add numbered steps — Claude follows ordered lists reliably"]),
        ("# heading doesn't match folder name",
         "The slash command name is derived from the folder — SKILL.md # heading must match.",
         ["Folder: .claude/skills/my-skill/",
          "SKILL.md first line must be: # /my-skill",
          "Not: # My Skill  or  # /My-Skill"]),
    ]

    top = 1.28
    for i, (prob, cause, fixes) in enumerate(issues):
        add_rect(sl, 0.3, top, 12.7, 0.3, DVIOLET)
        txhd = add_textbox(sl, 0.5, top + 0.03, 12.3, 0.24)
        txfhd = clear_first_para(txhd)
        para(txfhd, prob, 12, bold=True, colour=WHITE)
        top += 0.3

        add_rect(sl, 0.3, top, 12.7, 0.3, _rgb(0x5C, 0x20, 0x00))
        txca = add_textbox(sl, 0.5, top + 0.03, 12.3, 0.24)
        txfca = clear_first_para(txca)
        para(txfca, cause, 10, colour=_rgb(0xFF, 0xCC, 0x88), italic=True)
        top += 0.3

        add_rect(sl, 0.3, top, 12.7, 0.55, LAVEND)
        txfx = add_textbox(sl, 0.5, top + 0.05, 12.3, 0.46)
        txffx = clear_first_para(txfx)
        for j, fix in enumerate(fixes):
            sep = "   " if j > 0 else ""
            para(txffx, f"{sep}→  {fix}", 10, colour=DVIOLET, space_before=(0 if j == 0 else 2))
        top += 0.55 + 0.08

    add_rect(sl, 0.3, top, 12.7, 0.62, _rgb(0x1A, 0x5C, 0x2E))
    txb = add_textbox(sl, 0.5, top + 0.08, 12.3, 0.48)
    txf = clear_first_para(txb)
    para(txf, ("Quickest debug: open a new session → type / → if your skill name appears, it's working. "
               "If not, check git log on GitHub to confirm the push reached the right branch."),
         11, colour=_rgb(0x88, 0xFF, 0xBB))

    add_slide_number(sl, n, T)


def s15_best_practices(prs, n, T):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, INDIGO)
    add_rect(sl, 0, 0, 0.22, 7.5, VIOLET)

    txb = add_textbox(sl, 0.6, 0.25, 12.5, 0.7)
    txf = clear_first_para(txb)
    para(txf, "Best Practices & Summary", 30, bold=True, colour=WHITE)

    add_rect(sl, 0.5, 1.1, 5.9, 5.75, _rgb(0x2D, 0x27, 0x6B))
    txb2 = add_textbox(sl, 0.7, 1.2, 5.5, 0.45)
    txf2 = clear_first_para(txb2)
    para(txf2, "10 Rules for Great Skills", 16, bold=True, colour=VIOLET)
    txb3 = add_textbox(sl, 0.7, 1.7, 5.5, 5.0)
    txf3 = clear_first_para(txb3)
    for line in [
        "✓  One skill per folder, one SKILL.md per skill",
        "✓  Folder name = slash command name exactly",
        "✓  # heading must match folder name",
        "✓  Keep Instructions concrete and numbered",
        "✓  Name files and commands explicitly",
        "✓  Specify output format — Claude won't guess",
        "✓  Keep SKILL.md under 100 lines — use",
        "    companion .md files for long content",
        "✓  Version-control skills with your code",
        "✓  Start a fresh session after every push",
        "✓  Test by running the skill end-to-end",
    ]:
        para(txf3, line, 13, colour=WHITE, space_before=4)

    add_rect(sl, 6.7, 1.1, 6.3, 5.75, _rgb(0x2D, 0x27, 0x6B))
    txb4 = add_textbox(sl, 6.9, 1.2, 6.0, 0.45)
    txf4 = clear_first_para(txb4)
    para(txf4, "5-Step Quick Reference", 16, bold=True, colour=VIOLET)
    txb5 = add_textbox(sl, 6.9, 1.7, 6.0, 5.0)
    txf5 = clear_first_para(txb5)
    for line in [
        "1.  mkdir -p .claude/skills/my-skill",
        "    touch .claude/skills/my-skill/SKILL.md",
        "",
        "2.  Write SKILL.md:",
        "    # /my-skill",
        "    ## Instructions",
        "    When the user runs /my-skill: …",
        "",
        "3.  git add .claude/skills/my-skill/",
        "    git commit -m 'Add /my-skill skill'",
        "    git push -u origin main",
        "",
        "4.  Open Claude App → Code tab → New Session",
        "    Connect the repository",
        "",
        "5.  Type /my-skill in the chat → done",
    ]:
        para(txf5, line, 12, colour=WHITE, space_before=3)

    txb6 = add_textbox(sl, 0.5, 7.05, 12.5, 0.35)
    txf6 = clear_first_para(txb6)
    para(txf6, "Skills are Markdown files committed to your repo — no servers, no plugins, just push and go.",
         11, colour=GREY, align=PP_ALIGN.CENTER)

    add_slide_number(sl, n, T)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

TOTAL = 15

def build() -> Path:
    prs = new_prs()
    s01_title(prs,              1, TOTAL)
    s02_what_are_skills(prs,    2, TOTAL)
    s03_prerequisites(prs,      3, TOTAL)
    s04_folder(prs,             4, TOTAL)
    s05_skillmd_header(prs,     5, TOTAL)
    s06_skillmd_instructions(prs, 6, TOTAL)
    s07_supporting_files(prs,   7, TOTAL)
    s08_commit_push(prs,        8, TOTAL)
    s09_open_app(prs,           9, TOTAL)
    s10_invoke(prs,            10, TOTAL)
    s11_passing_args(prs,      11, TOTAL)
    s12_real_example_cq(prs,   12, TOTAL)
    s13_real_example_pptx(prs, 13, TOTAL)
    s14_debug(prs,             14, TOTAL)
    s15_best_practices(prs,    15, TOTAL)

    out = Path("/home/user/clinic/Claude_Skills_Guide_Detailed.pptx")
    prs.save(str(out))
    print(f"Saved: {out}  ({out.stat().st_size // 1024} KB)")
    return out


if __name__ == "__main__":
    build()
