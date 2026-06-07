"""Generate 'How to Use Claude Skills on Claude Code Web via Mobile App' — 10 slides."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette — deep purple / indigo
# ---------------------------------------------------------------------------
INDIGO  = RGBColor(0x1E, 0x1B, 0x4B)   # #1E1B4B  deep indigo primary
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED  violet accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LAVEND  = RGBColor(0xF5, 0xF3, 0xFF)   # #F5F3FF  light lavender bg
GREY    = RGBColor(0x6B, 0x72, 0x80)
ORANGE  = RGBColor(0xEA, 0x58, 0x0C)   # warning orange
LVIOLET = RGBColor(0xC4, 0xB5, 0xFD)   # light violet text on dark bg
DVIOLET = RGBColor(0x4C, 0x1D, 0x95)   # darker violet for card headers
GREEN   = RGBColor(0x05, 0x96, 0x69)   # tick green

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers  (verbatim pattern from make_pptx.py)
# ---------------------------------------------------------------------------

def _rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_solid(shape, colour: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, colour: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    fill_solid(shape, colour)
    shape.line.fill.background()
    return shape


def tf(shape):
    return shape.text_frame


def para(txf, text: str, size: int, bold=False, colour=WHITE,
         align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = txf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return p


def add_textbox(slide, left, top, width, height, word_wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.text_frame.word_wrap = word_wrap
    return txb


def clear_first_para(txb):
    """Remove the empty first paragraph added by default."""
    txf = txb.text_frame
    p = txf.paragraphs[0]
    p.clear()
    return txf


def add_slide_number(slide, num: int, total: int) -> None:
    txb = add_textbox(slide, 12.5, 7.1, 0.7, 0.3)
    txf = clear_first_para(txb)
    para(txf, f"{num}/{total}", 9, colour=GREY, align=PP_ALIGN.RIGHT)


def header_bar(slide, title: str, subtitle: str = "") -> None:
    add_rect(slide, 0, 0, 13.33, 1.1, INDIGO)
    txb = add_textbox(slide, 0.4, 0.08, 12.5, 0.55)
    txf = clear_first_para(txb)
    para(txf, title, 22, bold=True, colour=WHITE)
    if subtitle:
        txb2 = add_textbox(slide, 0.4, 0.62, 12.5, 0.4)
        txf2 = clear_first_para(txb2)
        para(txf2, subtitle, 13, colour=LVIOLET)


def bullet_box(slide, left, top, width, height, items: list[tuple[str, int, bool]],
               bg=None, pad=0.1):
    if bg:
        add_rect(slide, left, top, width, height, bg)
    txb = add_textbox(slide, left + pad, top + pad, width - pad * 2, height - pad * 2)
    txf = clear_first_para(txb)
    for text, size, bold in items:
        c = WHITE if bg in (INDIGO, VIOLET, DVIOLET) else INDIGO
        para(txf, text, size, bold=bold, colour=c, space_before=2)
    return txb


def card(slide, left, top, width, height, heading: str, lines: list[str],
         head_colour=VIOLET, body_colour=LAVEND):
    # heading strip
    h = 0.38
    add_rect(slide, left, top, width, h, head_colour)
    txh = add_textbox(slide, left + 0.1, top + 0.05, width - 0.2, h - 0.05)
    txf = clear_first_para(txh)
    para(txf, heading, 13, bold=True, colour=WHITE)
    # body
    add_rect(slide, left, top + h, width, height - h, body_colour)
    txb = add_textbox(slide, left + 0.12, top + h + 0.1,
                      width - 0.24, height - h - 0.15)
    txf2 = clear_first_para(txb)
    for line in lines:
        para(txf2, line, 11, colour=INDIGO, space_before=1)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs, num, total):
    """Title slide — deep indigo bg, violet left accent bar."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, INDIGO)
    add_rect(sl, 0, 0, 0.22, 7.5, VIOLET)

    # Main title
    txb = add_textbox(sl, 0.6, 1.5, 12.3, 1.2)
    txf = clear_first_para(txb)
    para(txf, "Custom Skills for Claude Code", 40, bold=True,
         colour=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txb2 = add_textbox(sl, 0.6, 2.8, 12.3, 0.65)
    txf2 = clear_first_para(txb2)
    para(txf2, "Build · Deploy · Activate on Claude Web via Mobile App", 22,
         colour=LVIOLET, align=PP_ALIGN.LEFT)

    # Violet divider line
    add_rect(sl, 0.6, 3.6, 8.0, 0.04, VIOLET)

    # Tag line
    txb3 = add_textbox(sl, 0.6, 3.8, 12.3, 0.5)
    txf3 = clear_first_para(txb3)
    para(txf3, "Step-by-step setup guide  ·  Claude Code Web  ·  Mobile App",
         14, colour=_rgb(0xA5, 0x92, 0xF9), align=PP_ALIGN.LEFT)

    add_slide_number(sl, num, total)


def slide_02_what_are_skills(prs, num, total):
    """What Are Claude Skills — two cards."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "What Are Claude Skills?", "Slash commands you define — Claude executes")

    card(sl, 0.3, 1.25, 6.2, 5.85, "What is a Skill?",
         ["A custom slash command  (e.g. /make-pptx)",
          "",
          "Defined by a SKILL.md file in your repo",
          "",
          "Tells Claude exactly what to do when invoked",
          "",
          "Lives in  .claude/skills/<skill-name>/SKILL.md",
          "",
          "Auto-activated when Claude Code reads your repo"])

    card(sl, 6.8, 1.25, 6.2, 5.85, "Why Use Skills?",
         ["Automate repetitive project tasks",
          "",
          "Share workflows with your team via Git",
          "",
          "No server or plugin install needed",
          "",
          "Works across Claude Code web + mobile + IDE",
          "",
          "Version-controlled alongside your code"])

    add_slide_number(sl, num, total)


def slide_03_prerequisites(prs, num, total):
    """Prerequisites — three checklist cards."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Prerequisites", "Three things you need before creating a skill")

    prereqs = [
        ("GitHub Repository",
         ["✓  Public or private repo on GitHub",
          "✓  You have push access to main / active branch",
          "✓  Skills must be committed to the branch",
          "     the session uses"]),
        ("Claude Mobile App",
         ["✓  Latest Claude app installed (iOS or Android)",
          "✓  Signed in to your Anthropic account",
          "✓  Navigate to the Code tab at bottom"]),
        ("Claude Code Session",
         ["✓  Repository connected in a web session",
          "✓  Session reads .claude/ on startup",
          "✓  New session re-reads repo on every start"]),
    ]

    for i, (heading, lines) in enumerate(prereqs):
        lft = 0.3 + i * 4.35
        card(sl, lft, 1.25, 4.1, 5.85, heading, lines,
             head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, num, total)


def slide_04_folder_structure(prs, num, total):
    """Step 1 — Create the Folder Structure."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 1 — Create the Folder Structure",
               "The folder name becomes the slash command name")

    # Large dark code box
    add_rect(sl, 0.3, 1.25, 7.8, 5.3, INDIGO)
    txb = add_textbox(sl, 0.55, 1.38, 7.3, 5.05)
    txf = clear_first_para(txb)
    para(txf, "Directory tree:", 12, bold=True, colour=LVIOLET)
    code_lines = [
        "your-repo/",
        "└── .claude/",
        "    └── skills/",
        "        ├── make-pptx/",
        "        │   ├── SKILL.md       ← required",
        "        │   └── pptxgenjs.md   ← optional reference",
        "        └── clinical-query/",
        "            └── SKILL.md       ← required",
    ]
    for line in code_lines:
        para(txf, line, 13, colour=_rgb(0xC4, 0xB5, 0xFD), space_before=3)

    # Right column: rules
    add_rect(sl, 8.4, 1.25, 4.6, 2.5, INDIGO)
    txb2 = add_textbox(sl, 8.6, 1.35, 4.2, 2.3)
    txf2 = clear_first_para(txb2)
    para(txf2, "Key Rule", 13, bold=True, colour=VIOLET)
    para(txf2, "Folder name  =  slash command name", 12, colour=WHITE, space_before=4)
    para(txf2, "", 8, colour=WHITE)
    para(txf2, "make-pptx/   →   /make-pptx", 12, colour=LVIOLET, space_before=2)
    para(txf2, "clinical-query/  →  /clinical-query", 12, colour=LVIOLET, space_before=2)

    # Create command box
    add_rect(sl, 8.4, 4.0, 4.6, 2.55, _rgb(0x2D, 0x27, 0x6B))
    txb3 = add_textbox(sl, 8.6, 4.1, 4.2, 2.35)
    txf3 = clear_first_para(txb3)
    para(txf3, "Create it:", 12, bold=True, colour=VIOLET)
    para(txf3, "mkdir -p .claude/skills/", 11, colour=LVIOLET, space_before=4)
    para(txf3, "    your-skill-name/", 11, colour=LVIOLET, space_before=1)
    para(txf3, "", 8, colour=WHITE)
    para(txf3, "Then add SKILL.md inside the folder.", 11, colour=WHITE, space_before=4)

    add_slide_number(sl, num, total)


def slide_05_write_skill_md(prs, num, total):
    """Step 2 — Write the SKILL.md."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 2 — Write SKILL.md",
               "The file Claude reads when your slash command is invoked")

    # Left: code template
    add_rect(sl, 0.3, 1.25, 7.3, 5.9, INDIGO)
    txb = add_textbox(sl, 0.5, 1.35, 7.0, 5.7)
    txf = clear_first_para(txb)
    para(txf, "SKILL.md template:", 12, bold=True, colour=LVIOLET)
    template_lines = [
        "# /your-skill-name",
        "",
        "One-line description of what this skill does.",
        "",
        "## Usage",
        "/your-skill-name [optional args]",
        "",
        "## What this skill does",
        "1. Step one",
        "2. Step two",
        "3. Step three",
        "",
        "## Instructions",
        "When the user invokes this skill:",
        "- Do X",
        "- Then do Y",
        "- Finally deliver Z",
        "",
        "## Output format",
        "Describe expected output here.",
    ]
    for line in template_lines:
        c = VIOLET if line.startswith("#") else (LVIOLET if line.startswith("-") or line[:1].isdigit() else WHITE)
        para(txf, line, 11, colour=c, space_before=1)

    # Right: tips card
    card(sl, 7.9, 1.25, 5.1, 5.9, "Tips for a Great SKILL.md",
         ["The # heading must match the folder name",
          "",
          "Instructions section is what Claude",
          "actually follows — make it clear",
          "",
          "Be specific: name files, commands,",
          "and expected outputs explicitly",
          "",
          "Reference other .md files in the same",
          "folder for long reference guides",
          "",
          "Keep Usage section short — one line",
          "with optional args is enough"],
         head_colour=DVIOLET, body_colour=WHITE)

    add_slide_number(sl, num, total)


def slide_06_commit_push(prs, num, total):
    """Step 3 — Commit & Push to GitHub."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 3 — Commit & Push to GitHub",
               "Skills only activate once they are in the repository Claude reads")

    # Commands dark box
    add_rect(sl, 0.3, 1.25, 12.7, 3.4, INDIGO)
    txb = add_textbox(sl, 0.55, 1.38, 12.2, 3.15)
    txf = clear_first_para(txb)
    para(txf, "# Stage the skill files", 12, colour=GREY, italic=True)
    para(txf, "git add .claude/skills/your-skill-name/", 14, colour=LVIOLET, space_before=2)
    para(txf, "", 8, colour=WHITE)
    para(txf, "# Commit with a clear message", 12, colour=GREY, italic=True, space_before=4)
    para(txf, 'git commit -m "Add /your-skill-name custom skill"', 14, colour=LVIOLET, space_before=2)
    para(txf, "", 8, colour=WHITE)
    para(txf, "# Push to the branch your session will use", 12, colour=GREY, italic=True, space_before=4)
    para(txf, "git push origin main", 14, colour=LVIOLET, space_before=2)

    # Warning box (orange)
    add_rect(sl, 0.3, 4.85, 12.7, 1.15, _rgb(0x7C, 0x28, 0x00))
    txb2 = add_textbox(sl, 0.5, 4.92, 12.3, 1.0)
    txf2 = clear_first_para(txb2)
    para(txf2, "Warning", 13, bold=True, colour=ORANGE)
    para(txf2, ("Skills on a feature branch only activate in sessions pointed at that branch.  "
                "Merge to main for universal access."),
         12, colour=WHITE, space_before=3)

    # Tip box
    add_rect(sl, 0.3, 6.2, 12.7, 0.92, _rgb(0x2D, 0x27, 0x6B))
    txb3 = add_textbox(sl, 0.5, 6.28, 12.3, 0.78)
    txf3 = clear_first_para(txb3)
    para(txf3, ("Tip:  Claude Code re-reads .claude/ from the repo on every new session start.  "
                "Skills pushed after session start will appear in the next session."),
         12, colour=LVIOLET)

    add_slide_number(sl, num, total)


def slide_07_open_mobile(prs, num, total):
    """Step 4 — Open Claude Code on Mobile App."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 4 — Open Claude Code on Mobile App",
               "Connect your repository so Claude can read your skills")

    steps = [
        ("1", "Open Claude App\n→ tap Code tab"),
        ("2", "Tap New Session\n(or existing session)"),
        ("3", "Connect repository\n→ select your repo"),
        ("4", "Claude clones repo\n& scans .claude/ folder"),
        ("5", "Skills registered\n→ slash commands active"),
    ]

    box_w = 2.2
    box_h = 2.4
    gap   = 0.35
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    start_l = (13.33 - total_w) / 2

    for i, (num_str, label) in enumerate(steps):
        lft = start_l + i * (box_w + gap)
        top = 1.4
        # Box bg
        add_rect(sl, lft, top, box_w, box_h, INDIGO)
        # Step number circle (simulated with small rect)
        add_rect(sl, lft + 0.05, top + 0.05, 0.45, 0.42, VIOLET)
        txn = add_textbox(sl, lft + 0.05, top + 0.06, 0.45, 0.38)
        txfn = clear_first_para(txn)
        para(txfn, num_str, 14, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        # Label
        txl = add_textbox(sl, lft + 0.1, top + 0.6, box_w - 0.2, 1.65)
        txfl = clear_first_para(txl)
        for part in label.split("\n"):
            para(txfl, part, 12, colour=WHITE, space_before=2)
        # Arrow (not after last step)
        if i < len(steps) - 1:
            arr_l = lft + box_w + 0.05
            txarr = add_textbox(sl, arr_l, top + box_h / 2 - 0.15, gap, 0.35)
            txfarr = clear_first_para(txarr)
            para(txfarr, "→", 16, bold=True, colour=VIOLET, align=PP_ALIGN.CENTER)

    # Bottom note
    add_rect(sl, 0.3, 4.1, 12.7, 1.05, _rgb(0x2D, 0x27, 0x6B))
    txb = add_textbox(sl, 0.5, 4.18, 12.3, 0.9)
    txf = clear_first_para(txb)
    para(txf, ("Each new session re-reads the repo.  "
               "Skills pushed after session start → start a new session to pick them up."),
         13, colour=LVIOLET)

    add_slide_number(sl, num, total)


def slide_08_use_skill(prs, num, total):
    """Step 5 — Activate & Use Your Skill."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Step 5 — Use Your Custom Skill",
               "Type / in the chat input to see and trigger registered skills")

    # Left: how to trigger
    card(sl, 0.3, 1.25, 5.8, 5.85, "How to Trigger",
         ["Type  /  in the chat input",
          "",
          "Autocomplete shows registered skills",
          "",
          "Select or type the full command:",
          "  /make-pptx",
          "  /clinical-query diabetic kidney disease",
          "",
          "Claude reads SKILL.md and executes",
          "the instructions you wrote",
          "",
          "Pass arguments after the command name",
          "just like a terminal command"],
         head_colour=DVIOLET, body_colour=WHITE)

    # Right: live example
    add_rect(sl, 6.4, 1.25, 6.6, 5.85, INDIGO)
    txb = add_textbox(sl, 6.6, 1.35, 6.2, 0.4)
    txf = clear_first_para(txb)
    para(txf, "Live example:", 12, bold=True, colour=LVIOLET)

    example_lines = [
        ("You: /clinical-query", VIOLET, True),
        ("     diabetic patient with kidney", WHITE, False),
        ("     complications", WHITE, False),
        ("", WHITE, False),
        ("Claude: Retrieving top-5 notes...", LVIOLET, False),
        ("", WHITE, False),
        ("[1] score=0.985  Nephrology", _rgb(0x86, 0xEF, 0xAC), False),
        ("    Chronic kidney disease in diabetic", WHITE, False),
        ("[2] score=0.702  Endocrinology", _rgb(0x86, 0xEF, 0xAC), False),
        ("    DKA in type 1 diabetic patient", WHITE, False),
        ("...", GREY, False),
        ("", WHITE, False),
        ("Clinical Summary: The top results", LVIOLET, False),
        ("highlight diabetic nephropathy...", LVIOLET, False),
    ]
    txb2 = add_textbox(sl, 6.6, 1.85, 6.2, 5.1)
    txf2 = clear_first_para(txb2)
    for text, colour, bold in example_lines:
        para(txf2, text, 11, colour=colour, bold=bold, space_before=1)

    add_slide_number(sl, num, total)


def slide_09_project_skills(prs, num, total):
    """Skills Built for This Project — two large cards."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Skills Built for This Project",
               "Two custom slash commands in .claude/skills/")

    # Card 1: /clinical-query
    add_rect(sl, 0.3, 1.25, 6.1, 0.45, DVIOLET)
    txh1 = add_textbox(sl, 0.4, 1.28, 5.9, 0.38)
    txf_h1 = clear_first_para(txh1)
    para(txf_h1, "/clinical-query", 16, bold=True, colour=WHITE)

    add_rect(sl, 0.3, 1.7, 6.1, 5.4, WHITE)
    txb1 = add_textbox(sl, 0.5, 1.8, 5.7, 5.2)
    txf1 = clear_first_para(txb1)
    para(txf1, "Trigger:", 12, bold=True, colour=DVIOLET)
    para(txf1, "Clinical scenario in natural language", 11, colour=INDIGO, space_before=1)
    para(txf1, "", 6, colour=INDIGO)
    para(txf1, "Actions:", 12, bold=True, colour=DVIOLET, space_before=4)
    for line in ["Runs retrieval against indexed notes",
                 "Prints ranked table with similarity scores",
                 "Generates LLM clinical summary",
                 "Reports score spread across top-k results"]:
        para(txf1, f"•  {line}", 11, colour=INDIGO, space_before=1)
    para(txf1, "", 6, colour=INDIGO)
    para(txf1, "Example:", 12, bold=True, colour=DVIOLET, space_before=4)
    para(txf1, "/clinical-query COPD with oxygen desaturation",
         11, colour=VIOLET, space_before=1)
    para(txf1, "", 6, colour=INDIGO)
    para(txf1, "File:", 12, bold=True, colour=DVIOLET, space_before=4)
    para(txf1, ".claude/skills/clinical-query/SKILL.md", 11, colour=GREY, space_before=1)

    # Card 2: /make-pptx
    add_rect(sl, 6.9, 1.25, 6.1, 0.45, DVIOLET)
    txh2 = add_textbox(sl, 7.0, 1.28, 5.9, 0.38)
    txf_h2 = clear_first_para(txh2)
    para(txf_h2, "/make-pptx", 16, bold=True, colour=WHITE)

    add_rect(sl, 6.9, 1.7, 6.1, 5.4, WHITE)
    txb2 = add_textbox(sl, 7.1, 1.8, 5.7, 5.2)
    txf2 = clear_first_para(txb2)
    para(txf2, "Trigger:", 12, bold=True, colour=DVIOLET)
    para(txf2, "/make-pptx  (no args needed)", 11, colour=INDIGO, space_before=1)
    para(txf2, "", 6, colour=INDIGO)
    para(txf2, "Actions:", 12, bold=True, colour=DVIOLET, space_before=4)
    for line in ["Runs scripts/make_pptx.py",
                 "QAs the generated output",
                 "Delivers the .pptx file",
                 "Commits and pushes to the repo"]:
        para(txf2, f"•  {line}", 11, colour=INDIGO, space_before=1)
    para(txf2, "", 6, colour=INDIGO)
    para(txf2, "Based on:", 12, bold=True, colour=DVIOLET, space_before=4)
    para(txf2, "Official anthropics/skills PPTX skill with",
         11, colour=INDIGO, space_before=1)
    para(txf2, "full design principles and QA workflow", 11, colour=INDIGO, space_before=1)
    para(txf2, "", 6, colour=INDIGO)
    para(txf2, "File:", 12, bold=True, colour=DVIOLET, space_before=4)
    para(txf2, ".claude/skills/make-pptx/SKILL.md", 11, colour=GREY, space_before=1)

    add_slide_number(sl, num, total)


def slide_10_troubleshooting(prs, num, total):
    """Troubleshooting & Best Practices."""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, LAVEND)
    header_bar(sl, "Troubleshooting & Best Practices",
               "Common issues and how to avoid them")

    # Left: issues table
    add_rect(sl, 0.3, 1.25, 6.5, 0.45, INDIGO)
    txh = add_textbox(sl, 0.4, 1.28, 6.3, 0.38)
    txf_h = clear_first_para(txh)
    para(txf_h, "Common Issues", 13, bold=True, colour=WHITE)

    issues = [
        ("Skill not in / list", "Pushed to wrong branch → merge to main"),
        ("Skill not showing",   "Start a new session (repo re-read on start)"),
        ("Command not found",   "Check folder name matches the # heading exactly"),
        ("Claude ignores steps", "Be more specific in the Instructions section"),
    ]

    for i, (problem, fix) in enumerate(issues):
        bg = WHITE if i % 2 == 0 else _rgb(0xEC, 0xE9, 0xFF)
        top_r = 1.7 + i * 0.65
        add_rect(sl, 0.3, top_r, 6.5, 0.65, bg)
        txp = add_textbox(sl, 0.4, top_r + 0.08, 2.4, 0.52)
        txfp = clear_first_para(txp)
        para(txfp, problem, 10, bold=True, colour=DVIOLET)
        txs = add_textbox(sl, 2.9, top_r + 0.08, 3.7, 0.52)
        txfs = clear_first_para(txs)
        para(txfs, fix, 10, colour=INDIGO)

    # Right: best practices
    add_rect(sl, 7.1, 1.25, 5.9, 0.45, DVIOLET)
    txbp = add_textbox(sl, 7.2, 1.28, 5.7, 0.38)
    txf_bp = clear_first_para(txbp)
    para(txf_bp, "Best Practices", 13, bold=True, colour=WHITE)

    add_rect(sl, 7.1, 1.7, 5.9, 4.9, WHITE)
    txbb = add_textbox(sl, 7.3, 1.8, 5.5, 4.7)
    txfb = clear_first_para(txbb)
    best = [
        "✓  One skill per folder, one SKILL.md per skill",
        "✓  Name the folder exactly what you want after /",
        "✓  Keep Instructions concrete — name files and commands",
        "✓  Store supporting .md guides in the same folder",
        "✓  Version-control skills alongside the code they automate",
        "✓  Test by starting a fresh session after pushing",
    ]
    for line in best:
        para(txfb, line, 12, colour=INDIGO, space_before=5)

    # Bottom summary bar
    add_rect(sl, 0.3, 6.4, 12.7, 0.75, INDIGO)
    txbs = add_textbox(sl, 0.5, 6.5, 12.3, 0.55)
    txfs = clear_first_para(txbs)
    para(txfs, ("Skills are just Markdown files committed to your repo.  "
                "No servers, no plugins — just push and go."),
         13, colour=LVIOLET, align=PP_ALIGN.CENTER)

    add_slide_number(sl, num, total)


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

TOTAL = 10


def build() -> Presentation:
    prs = new_prs()

    slide_01_title(prs,            1, TOTAL)
    slide_02_what_are_skills(prs,  2, TOTAL)
    slide_03_prerequisites(prs,    3, TOTAL)
    slide_04_folder_structure(prs, 4, TOTAL)
    slide_05_write_skill_md(prs,   5, TOTAL)
    slide_06_commit_push(prs,      6, TOTAL)
    slide_07_open_mobile(prs,      7, TOTAL)
    slide_08_use_skill(prs,        8, TOTAL)
    slide_09_project_skills(prs,   9, TOTAL)
    slide_10_troubleshooting(prs, 10, TOTAL)

    return prs


if __name__ == "__main__":
    out = Path("/home/user/clinic/Claude_Skills_Guide.pptx")
    prs = build()
    prs.save(str(out))
    print(f"Saved: {out}  ({out.stat().st_size // 1024} KB)")
